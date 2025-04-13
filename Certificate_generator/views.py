from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.models import User
from django.contrib.auth.decorators import login_required, user_passes_test
from django.contrib import messages
from django.contrib.auth import authenticate, login as auth_login
from django.contrib.auth.password_validation import validate_password
from django.core.exceptions import ValidationError
from django.core.files.storage import FileSystemStorage
from django.utils import timezone

from .models import Event
from .forms import CertificateForm, EventForm

from pptx import Presentation
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders

import csv
import os


def homepage(request):
    upcoming_events = Event.objects.all().order_by('-date')[:3]
    context = {'events': upcoming_events}
    return render(request, "homepage.html", context)


def signup(request):
    if request.method == 'POST':
        First_name = request.POST.get('First_name')
        Last_name = request.POST.get('Last_name')
        User_name = request.POST.get('User_name')
        email = request.POST.get('email')
        Password1 = request.POST.get('Password1')
        Password2 = request.POST.get('Password2')

        if User.objects.filter(username=User_name).exists():
            message = "Username already taken. Please choose a different one."
            return render(request, "signup.html", {'message': message})

        if Password1 != Password2:
            message = "Passwords do not match. Please try again."
            return render(request, "signup.html", {'message': message})

        try:
            validate_password(Password1)
        except ValidationError as e:
            message = ', '.join(e.messages)
            return render(request, "signup.html", {'message': message})

        my_user = User.objects.create_user(
            username=User_name,
            email=email,
            password=Password1,
            first_name=First_name,
            last_name=Last_name
        )
        my_user.save()
        return redirect('homepage')

    return render(request, "signup.html")


def login(request):
    if request.method == 'POST':
        username = request.POST.get('User_name')
        password = request.POST.get('Password')
        user = authenticate(request, username=username, password=password)
        if user is not None:
            auth_login(request, user)
            return redirect('homepage')
        else:
            message = "Invalid username or password. Please try again."
            return render(request, "signup.html", {'message': message})
    return render(request, "signup.html")


def generate_certificate(request):
    sent_certificates = []
    if request.method == "POST":
        form = CertificateForm(request.POST, request.FILES)
        if form.is_valid():
            template = form.cleaned_data['template']
            student_details_file = form.cleaned_data['csvfile']
            sender_email = form.cleaned_data['sender_email']
            sender_password = form.cleaned_data['sender_password']

            fs = FileSystemStorage()
            filename = fs.save(student_details_file.name, student_details_file)

            try:
                detail = []
                with open(fs.path(filename), 'r', encoding='utf-8') as csvfile:
                    csvreader = csv.reader(csvfile)
                    for row in csvreader:
                        detail.append([cell.strip() for cell in row])

                keys = detail[0]
                student_dict = {}
                for row in detail[1:]:
                    student_dict[row[0]] = dict(zip(keys[1:], row[1:]))

                for key, values in student_dict.items():
                    prs = Presentation(template)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text_frame = shape.text_frame
                                for paragraph in text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        for old_text, new_text in values.items():
                                            if old_text.upper() in run.text:
                                                run.text = run.text.replace(old_text.upper(), new_text)
                    certificate_file = f"{values['name']}_Certificate.pptx"
                    prs.save(certificate_file)
                    send_certificate(sender_email, sender_password, values['email'], values['name'], certificate_file)
                    sent_certificates.append({'name': values['name'], 'email': values['email']})
                    messages.success(request, f"Email sent successfully to {values['name']} ({values['email']})")
            except Exception as e:
                print("Error processing CSV data:", e)
            finally:
                fs.delete(filename)
        return redirect('mail_status')
    else:
        form = CertificateForm()
    return render(request, "generatecerti.html", {"form": form, "sent_certificates": sent_certificates})


def send_certificate(sender_email, sender_password, recipient_email, recipient_name, certificate_file):
    smtp_server = "smtp.gmail.com"
    smtp_port = 587
    msg = MIMEMultipart()
    msg['From'] = sender_email
    msg['To'] = recipient_email
    msg['Subject'] = "Certificate"
    with open(certificate_file, "rb") as file:
        attachment = MIMEBase("application", "vnd.openxmlformats-officedocument.presentationml.presentation")
        attachment.set_payload(file.read())
        encoders.encode_base64(attachment)
        attachment.add_header('Content-Disposition', f'attachment; filename="{os.path.basename(certificate_file)}"')
        msg.attach(attachment)
    try:
        server = smtplib.SMTP(smtp_server, smtp_port)
        server.starttls()
        server.login(sender_email, sender_password)
        server.sendmail(sender_email, recipient_email, msg.as_string())
        print(f"Certificate sent to {recipient_email}")
        server.quit()
        if os.path.exists(certificate_file):
            os.remove(certificate_file)
            print(f"Certificate file {certificate_file} deleted.")
        else:
            print(f"Certificate file {certificate_file} does not exist.")
    except Exception as e:
        print("Error:", e)


def mail_status(request):
    return render(request, "mail_status.html")


def create_event(request):
    if request.method == 'POST':
        form = EventForm(request.POST, request.FILES)
        if form.is_valid():
            form.save()
            return redirect('event_list')
    else:
        form = EventForm()
    return render(request, 'create_event.html', {'form': form})


def event_list(request):
    events = Event.objects.all().order_by('date')
    now = timezone.now()
    context = {'events': events, 'now': now}
    return render(request, 'events.html', context)


def participate_event(request, event_id):
    event = get_object_or_404(Event, pk=event_id)
    user = request.user

    if user in event.joined_by.all():
        messages.error(request, "You have already participated in this event.")
        return redirect('event_list')

    if event.current_participants >= event.participants:
        messages.error(request, "This event is full.")
        return redirect('event_list')

    event.joined_by.add(user)
    event.current_participants += 1
    event.save()

    messages.success(request, "You have successfully participated in this event!")
    return redirect('event_list')


def event_participants(request, event_id):
    event = get_object_or_404(Event, pk=event_id)
    participants = event.joined_by.all()
    return render(request, 'event_participants.html', {'event': event, 'participants': participants})


@login_required
@user_passes_test(lambda u: u.is_staff)
def delete_event(request, event_id):
    event = get_object_or_404(Event, pk=event_id)
    if request.method == 'POST':
        event.delete()
    return redirect('event_list')


@login_required
def user_profile_view(request):
    user = request.user
    return render(request, 'user_profile_view.html', {'user': user})
