from rest_framework.views import APIView
from rest_framework.parsers import MultiPartParser, FormParser
from rest_framework.response import Response
from rest_framework import status
from Certificate_generator.api.serializers import CertificateRequestSerializer
from Certificate_generator.views import send_certificate # Reuse your existing function
from django.utils.decorators import method_decorator
from django.views.decorators.csrf import csrf_exempt

from pptx import Presentation
import csv
import os
from django.core.files.storage import FileSystemStorage




@method_decorator(csrf_exempt, name='dispatch')
class CertificateGenerateAPI(APIView):

    parser_classes = (MultiPartParser, FormParser)

    def post(self, request, format=None):
        serializer = CertificateRequestSerializer(data=request.data)
        if serializer.is_valid():
            template = request.FILES['template']
            csvfile = request.FILES['csvfile']
            sender_email = serializer.validated_data['sender_email']
            sender_password = serializer.validated_data['sender_password']

            fs = FileSystemStorage()
            template_path = fs.save(template.name, template)
            csv_path = fs.save(csvfile.name, csvfile)

            try:
                with open(fs.path(csv_path), 'r', encoding='utf-8') as file:
                    csvreader = csv.reader(file)
                    rows = [row for row in csvreader]
                
                headers = rows[0]
                for row in rows[1:]:
                    data = dict(zip(headers[1:], row[1:]))
                    prs = Presentation(fs.path(template_path))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    for run in para.runs:
                                        for old_text, new_text in data.items():
                                            if old_text.upper() in run.text:
                                                run.text = run.text.replace(old_text.upper(), new_text)

                    output_file = f"{data['name']}_Certificate.pptx"
                    prs.save(output_file)

                    send_certificate(
                        sender_email, sender_password,
                        data['email'], data['name'],
                        output_file
                    )

            except Exception as e:
                return Response({"error": str(e)}, status=500)
            finally:
                fs.delete(template_path)
                fs.delete(csv_path)

            return Response({"message": "Certificates sent successfully."}, status=200)
        return Response(serializer.errors, status=400)
